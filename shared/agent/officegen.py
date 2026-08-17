"""Génération de fichiers bureautiques pour VindIA (Word, Excel, PowerPoint, PDF).

VindIA produit du texte structuré (markdown) entre [[FICHIER:nom.ext]] ; ce module le
convertit en VRAI binaire. Aucun stockage : construit en mémoire, renvoie les octets.

Mise en page gérée :
  - Titres « # / ## / ### » (colorés à la charte), paragraphes
  - Puces « - » et listes numérotées « 1. »
  - GRAS « **texte** » conservé (Word et PDF)
  - TABLEAUX markdown «| col | col |» avec ligne de séparation «|---|---|» (Word et PDF)
  - .xlsx = CSV → feuille ; .pptx = diapos séparées par «---»
"""

from __future__ import annotations

import csv
import io
import re
from pathlib import Path

_DEJAVU = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
_DEJAVU_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

# Couleur d'accent (charte VindIA, indigo) pour les titres.
_ACCENT = (79, 70, 229)

# Images : insérées via « ![alt](nom) » si le fichier existe sous base_dir.
_IMG_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}

# Transparence IA — AI Act art. 50(2) : tout contenu généré doit être identifiable
# comme tel. Mention lisible par un humain + métadonnées lisibles par une machine.
AI_NOTICE = "Document généré par intelligence artificielle (VindIA) — à vérifier avant usage."
AI_META = "AI-generated content (VindIA) — EU AI Act art. 50"

# Provenance XMP pour les PDF. L'AI Act est neutre technologiquement : il exige un
# marquage lisible par machine, et cite les « métadonnées d'identification » — ce que
# fournit XMP, sans imposer la signature cryptographique d'un C2PA complet.
# Les champs suivent les vocabulaires standards (Dublin Core, XMP Rights, PDF/A extension)
# pour rester exploitables par un outil d'analyse tiers.
def _xmp_provenance() -> str:
    # NB : ne PAS inclure l'enveloppe <?xpacket…?> — fpdf2 l'ajoute lui-même et
    # refuse (ValueError) un XMP déjà enveloppé.
    return (
        '<x:xmpmeta xmlns:x="adobe:ns:meta/">'
        '<rdf:RDF xmlns:rdf="http://www.w3.org/1999/02/22-rdf-syntax-ns#">'
        '<rdf:Description rdf:about=""'
        ' xmlns:dc="http://purl.org/dc/elements/1.1/"'
        ' xmlns:xmpRights="http://ns.adobe.com/xap/1.0/rights/"'
        ' xmlns:xmp="http://ns.adobe.com/xap/1.0/">'
        '<dc:creator><rdf:Seq><rdf:li>VindIA (IA)</rdf:li></rdf:Seq></dc:creator>'
        '<dc:type><rdf:Bag><rdf:li>AI-generated</rdf:li></rdf:Bag></dc:type>'
        f'<dc:description><rdf:Alt><rdf:li xml:lang="x-default">{AI_META}</rdf:li></rdf:Alt></dc:description>'
        '<xmp:CreatorTool>VindIA — generative AI assistant</xmp:CreatorTool>'
        '<xmpRights:Marked>True</xmpRights:Marked>'
        '<xmpRights:UsageTerms><rdf:Alt><rdf:li xml:lang="x-default">'
        'Contenu généré par intelligence artificielle. Vérifier avant usage.'
        '</rdf:li></rdf:Alt></xmpRights:UsageTerms>'
        '</rdf:Description></rdf:RDF></x:xmpmeta>'
    )


# En-tête ajouté aux fichiers TEXTE générés (md, txt, csv…) : marquage lisible par
# un humain ET repérable par une machine (ligne préfixée, format stable).
TEXT_MARK = f"<!-- {AI_META} -->"


def _tag_office_meta(props) -> None:
    """Marque les propriétés d'un document Office comme contenu généré par IA."""
    try:
        props.comments = AI_META
        props.category = "AI-generated"
        props.author = "VindIA (IA)"
    except Exception:
        pass


def _safe_image(base_dir, rel: str):
    """Chemin absolu d'une image SOUS base_dir (anti path-traversal), ou None."""
    if not base_dir:
        return None
    try:
        base = Path(base_dir).resolve()
        target = (base / rel.lstrip("/")).resolve()
    except Exception:
        return None
    if (target == base or base in target.parents) and target.is_file() and target.suffix.lower() in _IMG_EXT:
        return target
    return None

OFFICE_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
}


# --------------------------------------------------------------------------- #
#  Analyse du markdown (segments gras, blocs, tableaux)
# --------------------------------------------------------------------------- #

def _inline_segments(text: str):
    """Découpe une ligne en segments (texte, gras) en gérant **gras** et `code`."""
    text = re.sub(r"`(.+?)`", r"\1", text)
    segs = []
    for i, part in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if part != "":
            segs.append((part, i % 2 == 1))  # parties impaires = entre ** = gras
    return segs or [(text, False)]


def _strip_inline(text: str) -> str:
    """Texte plat : retire le gras/italique/`code` (pour titres, cellules de tableau)."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1", text)
    return text


def _is_table_sep(line: str) -> bool:
    """Ligne de séparation d'un tableau markdown, ex « |---|:--:|--- »."""
    s = line.strip()
    return "|" in s and "-" in s and bool(re.match(r"^\|?[\s:|-]+\|?$", s))


def _table_row(line: str):
    return [c.strip() for c in line.strip().strip("|").split("|")]


def _blocks(content: str):
    """Itère les blocs : ('h1'|'h2'|'h3'|'bullet'|'num'|'para'|'table'|'blank', data)."""
    lines = content.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        # Tableau : une ligne avec « | » suivie d'une ligne de séparation.
        if "|" in line and line.strip() and i + 1 < len(lines) and _is_table_sep(lines[i + 1]):
            rows = [_table_row(line)]
            i += 2
            while i < len(lines) and "|" in lines[i] and lines[i].strip():
                rows.append(_table_row(lines[i]))
                i += 1
            yield ("table", rows)
            continue
        img = re.match(r"^\s*!\[([^\]]*)\]\(([^)]+)\)\s*$", line)
        if img:
            yield ("image", (img.group(1), img.group(2).strip()))
            i += 1
            continue
        if not line.strip():
            yield ("blank", None)
        elif line.startswith("### "):
            yield ("h3", line[4:].strip())
        elif line.startswith("## "):
            yield ("h2", line[3:].strip())
        elif line.startswith("# "):
            yield ("h1", line[2:].strip())
        elif re.match(r"^\s*[-*]\s+", line):
            yield ("bullet", re.sub(r"^\s*[-*]\s+", "", line))
        elif re.match(r"^\s*\d+[.)]\s+", line):
            yield ("num", re.sub(r"^\s*\d+[.)]\s+", "", line))
        else:
            yield ("para", line)
        i += 1


# --------------------------------------------------------------------------- #
#  Word
# --------------------------------------------------------------------------- #

def _build_docx(content: str, base_dir=None) -> bytes:
    from docx import Document
    from docx.shared import Inches, RGBColor

    doc = Document()
    accent = RGBColor(*_ACCENT)
    for kind, data in _blocks(content):
        if kind == "blank":
            continue
        if kind in ("h1", "h2", "h3"):
            h = doc.add_heading(level=int(kind[1]))
            run = h.add_run(_strip_inline(data))
            run.font.color.rgb = accent
        elif kind == "bullet":
            _docx_runs(doc.add_paragraph(style="List Bullet"), data)
        elif kind == "num":
            _docx_runs(doc.add_paragraph(style="List Number"), data)
        elif kind == "table":
            _docx_table(doc, data)
        elif kind == "image":
            alt, rel = data
            img = _safe_image(base_dir, rel)
            if img:
                try:
                    doc.add_picture(str(img), width=Inches(4.0))
                except Exception:
                    doc.add_paragraph(f"[image : {alt or rel}]")
            else:
                doc.add_paragraph(f"[image introuvable : {alt or rel}]")
        else:
            _docx_runs(doc.add_paragraph(), data)
    # Transparence IA (AI Act art. 50) : mention visible + métadonnées machine.
    note = doc.add_paragraph()
    run = note.add_run(AI_NOTICE)
    run.italic = True
    try:
        from docx.shared import Pt
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x77, 0x77, 0x77)
    except Exception:
        pass
    _tag_office_meta(doc.core_properties)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _docx_runs(paragraph, text: str) -> None:
    for seg, bold in _inline_segments(text):
        run = paragraph.add_run(seg)
        run.bold = bold


def _docx_table(doc, rows) -> None:
    ncols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=ncols)
    try:
        table.style = "Light Grid Accent 1"
    except Exception:
        table.style = "Table Grid"
    for ri, row in enumerate(rows):
        for ci in range(ncols):
            val = row[ci] if ci < len(row) else ""
            cell = table.rows[ri].cells[ci]
            run = cell.paragraphs[0].add_run(_strip_inline(val))
            if ri == 0:
                run.bold = True


# --------------------------------------------------------------------------- #
#  Excel / PowerPoint
# --------------------------------------------------------------------------- #

def _sniff_rows(content: str):
    text = content.strip("\n")
    sample = text[:2000]
    delim = ","
    try:
        delim = csv.Sniffer().sniff(sample, delimiters=",;\t").delimiter
    except Exception:
        for cand in (";", "\t", ","):
            if cand in sample:
                delim = cand
                break
    return list(csv.reader(io.StringIO(text), delimiter=delim))


def _build_xlsx(content: str, base_dir=None) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = "Feuille1"
    rows = _sniff_rows(content)
    # Mise en forme : un tableur livrable se lit d'un coup d'œil — en-tête contrasté,
    # filtres, première ligne figée, nombres alignés à droite.
    entete_fond = PatternFill("solid", fgColor="4F46E5")
    filet = Side(style="thin", color="D9D9E3")
    bordure = Border(left=filet, right=filet, top=filet, bottom=filet)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row, start=1):
            valeur, texte = val, str(val).strip()
            # Un nombre écrit en texte ne se totalise pas : on convertit quand c'est
            # sans ambiguïté (virgule décimale française comprise).
            if r > 1 and texte:
                essai = texte.replace("\u202f", "").replace(" ", "").replace(",", ".")
                if essai.replace(".", "", 1).replace("-", "", 1).isdigit():
                    try:
                        valeur = float(essai) if "." in essai else int(essai)
                    except ValueError:
                        valeur = val
            cell = ws.cell(row=r, column=c, value=valeur)
            cell.border = bordure
            if r == 1:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = entete_fond
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            elif isinstance(valeur, (int, float)):
                cell.alignment = Alignment(horizontal="right")
    if len(rows) > 1:
        ws.freeze_panes = "A2"
        ws.auto_filter.ref = f"A1:{get_column_letter(max(len(r) for r in rows))}{len(rows)}"
    if rows:
        for c in range(1, max(len(r) for r in rows) + 1):
            width = max((len(str(row[c - 1])) for row in rows if len(row) >= c), default=10)
            ws.column_dimensions[ws.cell(row=1, column=c).column_letter].width = min(max(width + 2, 10), 60)
    # Transparence IA (AI Act art. 50).
    last = ws.max_row + 2
    cell = ws.cell(row=last, column=1, value=AI_NOTICE)
    cell.font = Font(italic=True, size=8, color="777777")
    try:
        wb.properties.creator = "VindIA (IA)"
        wb.properties.description = AI_META
        wb.properties.category = "AI-generated"
    except Exception:
        pass
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_pptx(content: str, base_dir=None) -> bytes:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    layout = prs.slide_layouts[1]
    chunks = re.split(r"(?m)^\s*---\s*$", content)
    chunks = [c for c in chunks if c.strip()] or [content]
    for chunk in chunks:
        lines = [l.rstrip() for l in chunk.splitlines() if l.strip()]
        if not lines:
            continue
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = _strip_inline(re.sub(r"^#+\s*", "", lines[0]))
        body = slide.placeholders[1].text_frame
        body.clear()
        first = True
        for line in lines[1:]:
            txt = _strip_inline(re.sub(r"^\s*[-*]\s+", "", line))
            p = body.paragraphs[0] if first else body.add_paragraph()
            p.text = txt
            p.font.size = Pt(18)
            first = False
    _tag_office_meta(prs.core_properties)   # transparence IA (AI Act art. 50)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
#  PDF
# --------------------------------------------------------------------------- #

def _build_pdf(content: str, base_dir=None) -> bytes:
    from fpdf import FPDF

    pdf = FPDF(format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.add_font("DejaVu", "", _DEJAVU)
    pdf.add_font("DejaVu", "B", _DEJAVU_BOLD)
    pdf.set_font("DejaVu", "", 11)

    for kind, data in _blocks(content):
        if kind == "blank":
            pdf.ln(3)
        elif kind in ("h1", "h2", "h3"):
            size = {"h1": 18, "h2": 15, "h3": 13}[kind]
            pdf.set_text_color(*_ACCENT)
            pdf.set_font("DejaVu", "B", size)
            pdf.multi_cell(pdf.epw, size * 0.5, _strip_inline(data), new_x="LMARGIN", new_y="NEXT")
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1)
        elif kind == "bullet":
            _pdf_line(pdf, data, prefix="  •  ")
        elif kind == "num":
            _pdf_line(pdf, data, prefix="  ")
        elif kind == "table":
            _pdf_table(pdf, data)
        elif kind == "image":
            alt, rel = data
            img = _safe_image(base_dir, rel)
            if img:
                try:
                    pdf.image(str(img), w=min(90, pdf.epw))
                    pdf.ln(3)
                except Exception:
                    _pdf_line(pdf, f"[image : {alt or rel}]")
            else:
                _pdf_line(pdf, f"[image introuvable : {alt or rel}]")
        else:
            _pdf_line(pdf, data)
    # Transparence IA (AI Act art. 50) : mention visible + métadonnées machine.
    pdf.ln(4)
    pdf.set_font("DejaVu", "", 8)
    pdf.set_text_color(120, 120, 120)
    pdf.multi_cell(pdf.epw, 4, AI_NOTICE, new_x="LMARGIN", new_y="NEXT")
    pdf.set_text_color(0, 0, 0)
    try:
        pdf.set_author("VindIA (IA)")
        pdf.set_creator("VindIA — AI-generated")
        pdf.set_subject(AI_META)
        pdf.set_keywords("AI-generated, VindIA, EU AI Act art.50")
        pdf.set_producer("VindIA (generative AI)")
    except Exception:
        pass
    try:
        pdf.set_xmp_metadata(_xmp_provenance())   # provenance normalisée (XMP)
    except Exception as exc:  # ne pas masquer : un XMP invalide doit se voir
        print(f"[VindIA] XMP non appliqué : {exc}")
    return bytes(pdf.output())


def _pdf_line(pdf, text: str, prefix: str = "") -> None:
    """Écrit une ligne avec gras inline (write gère le retour à la ligne)."""
    pdf.set_x(pdf.l_margin)
    if prefix:
        pdf.set_font("DejaVu", "", 11)
        pdf.write(6, prefix)
    for seg, bold in _inline_segments(text):
        pdf.set_font("DejaVu", "B" if bold else "", 11)
        pdf.write(6, seg)
    pdf.ln(8)


def _pdf_table(pdf, rows) -> None:
    pdf.set_font("DejaVu", "", 10)
    with pdf.table() as table:
        for row in rows:
            tr = table.row()
            for cell in row:
                tr.cell(_strip_inline(cell))
    pdf.set_font("DejaVu", "", 11)
    pdf.ln(2)


_BUILDERS = {
    "docx": _build_docx,
    "xlsx": _build_xlsx,
    "pptx": _build_pptx,
    "pdf": _build_pdf,
}


def build_file(name: str, content: str, base_dir=None) -> tuple[bytes, str]:
    """Construit le fichier binaire. `base_dir` permet d'insérer des images locales
    (« ![alt](nom) ») trouvées sous ce dossier. Lève ValueError si l'extension n'est
    pas gérée."""
    ext = (name.rsplit(".", 1)[-1] if "." in name else "").lower()
    if ext not in _BUILDERS:
        raise ValueError(f"format non supporte: {ext}")
    return _BUILDERS[ext](content or "", base_dir), OFFICE_TYPES[ext]
