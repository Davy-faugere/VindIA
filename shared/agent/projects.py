"""Projets persistants par utilisateur : VindIA apprend des fichiers déposés.

But (Lot 2a) : l'utilisateur dépose des documents dans un PROJET ; VindIA les
ingère (extraction texte), les range dans un espace PRIVÉ persistant, et peut s'y
référer plus tard pour « suivre le projet ». Deux garanties de fond :

  - ISOLATION STRICTE par membre. Tout est rangé sous `<base>/<member_id>/...`.
    `member_id`, `project_id` et noms de fichiers sont assainis (anti
    path-traversal) : un membre ne peut JAMAIS lire/écrire l'espace d'un autre,
    même avec une entrée malveillante. C'est le code qui le garantit, pas une
    convention d'appel.
  - PERSISTANCE sur DISQUE (indépendant de MariaDB) : survit aux sessions et aux
    redémarrages du service.

Style maison : ce module n'importe AUCUNE dépendance tierce au chargement.
L'extraction des formats binaires (docx/xlsx/pptx/pdf) importe les libs
PARESSEUSEMENT, au premier appel seulement → la CI stdlib reste 0-dépendance.
"""

from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, List, Optional

from .tools import html_to_text

# --------------------------------------------------------------------------- #
#  Assainissement (le garde-fou d'isolation)
# --------------------------------------------------------------------------- #

# member_id VindIA = UUID CHAR(36) → uniquement hex + tirets. Tout le reste est
# rejeté : pas de séparateur de chemin, pas de « .. », pas d'absolu.
_MEMBER_RE = re.compile(r"^[0-9a-fA-F-]{1,36}$")
_SLUG_RE = re.compile(r"[^a-z0-9]+")


def _safe_member(member_id: str) -> str:
    if not member_id or not _MEMBER_RE.match(member_id):
        raise ValueError("member_id invalide")
    return member_id


def slugify(name: str, *, fallback: str = "projet") -> str:
    """Nom lisible -> slug sûr pour un nom de dossier (a-z0-9-, borné)."""
    s = _SLUG_RE.sub("-", (name or "").strip().lower()).strip("-")
    return (s or fallback)[:48]


def safe_filename(name: str) -> str:
    """Réduit un nom de fichier à son basename assaini (jamais de traversée)."""
    base = Path(name or "").name  # retire tout composant de chemin
    base = base.replace("\x00", "")
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", base).strip("._") or "fichier"
    return cleaned[:120]


# --------------------------------------------------------------------------- #
#  Extraction de texte (ingestion)
# --------------------------------------------------------------------------- #

# Formats décodés en texte pur (stdlib, testables offline).
_TEXT_EXT = {".txt", ".md", ".markdown", ".csv", ".log", ".json", ".yaml", ".yml", ".tsv"}
_HTML_EXT = {".html", ".htm"}


class ExtractionError(RuntimeError):
    """Extraction impossible (format non supporté ou contenu illisible)."""


def extract_text(filename: str, data: bytes) -> str:
    """Extrait le texte d'un fichier d'après son extension. Lève si non géré.

    Texte/HTML : stdlib pur. docx/xlsx/pptx : libs importées paresseusement (déjà
    installées pour filegen). pdf : best-effort si une lib de lecture est présente.
    """
    ext = Path(filename or "").suffix.lower()
    if ext in _TEXT_EXT:
        return data.decode("utf-8", errors="replace").strip()
    if ext in _HTML_EXT:
        return html_to_text(data.decode("utf-8", errors="replace"))
    if ext == ".docx":
        return _extract_docx(data)
    if ext == ".xlsx":
        return _extract_xlsx(data)
    if ext == ".pptx":
        return _extract_pptx(data)
    if ext == ".pdf":
        return _extract_pdf(data)
    raise ExtractionError(f"format non supporté : {ext or '(sans extension)'}")


def _extract_docx(data: bytes) -> str:  # pragma: no cover - dépend de python-docx
    import io

    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:  # inclut aussi le contenu des tableaux
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_xlsx(data: bytes) -> str:  # pragma: no cover - dépend d'openpyxl
    import io

    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: List[str] = []
    for ws in wb.worksheets:
        lines.append(f"# Feuille : {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    return "\n".join(lines).strip()


def _extract_pptx(data: bytes) -> str:  # pragma: no cover - dépend de python-pptx
    import io

    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"# Diapo {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
    return "\n".join(lines).strip()


def _extract_pdf(data: bytes) -> str:  # pragma: no cover - lib de lecture optionnelle
    import io

    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader  # type: ignore
        except ImportError as exc:
            raise ExtractionError(
                "lecture PDF indisponible : installer `pypdf` dans le venv."
            ) from exc
    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages).strip()


# --------------------------------------------------------------------------- #
#  Magasin de projets
# --------------------------------------------------------------------------- #

@dataclass(frozen=True)
class Document:
    filename: str
    chars: int
    added_at: str


@dataclass(frozen=True)
class Project:
    project_id: str
    name: str
    created_at: str
    documents: List[Document]

    def as_dict(self) -> dict:
        return {
            "project_id": self.project_id,
            "name": self.name,
            "created_at": self.created_at,
            "documents": [d.__dict__ for d in self.documents],
        }


# Horloge injectable : par défaut UTC ISO ; en test, une horloge fixe → déterminisme.
Clock = Callable[[], str]


def _default_clock() -> str:  # pragma: no cover - dépend de l'heure réelle
    import datetime

    return datetime.datetime.now(datetime.timezone.utc).replace(microsecond=0).isoformat()


class ProjectStore:
    """CRUD de projets isolés par membre, persistés sur disque.

    Arborescence : `<base>/<member_id>/<project_id>/meta.json` + `docs/<fichier>`.
    Toute entrée traverse l'assainissement → impossible de sortir de l'espace du
    membre (`..`, chemins absolus, séparateurs sont neutralisés ou rejetés).
    """

    # Bornes anti-abus : contexte injecté au LLM (system prompt) plafonné.
    CONTEXT_BUDGET = 6000

    def __init__(self, base_dir: str, *, clock: Optional[Clock] = None) -> None:
        self._base = Path(base_dir)
        self._clock = clock or _default_clock

    # -- chemins (toujours via l'assainissement) ----------------------------- #
    def _member_dir(self, member_id: str) -> Path:
        return self._base / _safe_member(member_id)

    def _project_dir(self, member_id: str, project_id: str) -> Path:
        pid = slugify(project_id)
        d = (self._member_dir(member_id) / pid).resolve()
        # Défense en profondeur : le dossier résolu DOIT rester sous le membre.
        member_root = self._member_dir(member_id).resolve()
        if member_root not in d.parents and d != member_root:
            raise ValueError("chemin de projet hors de l'espace membre")
        return d

    # -- opérations ---------------------------------------------------------- #
    def create_project(self, member_id: str, name: str) -> Project:
        """Crée un projet (slug du nom, suffixe -2/-3… en cas de collision)."""
        base_slug = slugify(name)
        member_dir = self._member_dir(member_id)
        pid = base_slug
        i = 2
        while (member_dir / pid).exists():
            pid = f"{base_slug}-{i}"
            i += 1
        pdir = member_dir / pid
        (pdir / "docs").mkdir(parents=True, exist_ok=True)
        proj = Project(pid, (name or pid).strip()[:120], self._clock(), [])
        self._write_meta(member_id, proj)
        return proj

    def list_projects(self, member_id: str) -> List[Project]:
        member_dir = self._member_dir(member_id)
        if not member_dir.exists():
            return []
        out: List[Project] = []
        for child in sorted(member_dir.iterdir()):
            meta = child / "meta.json"
            if meta.is_file():
                out.append(self._read_meta(meta))
        return out

    def get_project(self, member_id: str, project_id: str) -> Optional[Project]:
        meta = self._project_dir(member_id, project_id) / "meta.json"
        return self._read_meta(meta) if meta.is_file() else None

    def add_document(self, member_id: str, project_id: str, filename: str, text: str) -> Document:
        """Range un texte ingéré dans le projet et met à jour ses métadonnées."""
        proj = self.get_project(member_id, project_id)
        if proj is None:
            raise ValueError("projet inconnu")
        fname = safe_filename(filename)
        pdir = self._project_dir(member_id, project_id)
        (pdir / "docs" / f"{fname}.txt").write_text(text, encoding="utf-8")
        doc = Document(fname, len(text), self._clock())
        # Remplace un éventuel doc de même nom (ré-upload), sinon ajoute.
        docs = [d for d in proj.documents if d.filename != fname] + [doc]
        self._write_meta(member_id, Project(proj.project_id, proj.name, proj.created_at, docs))
        return doc

    def read_document(self, member_id: str, project_id: str, filename: str) -> str:
        path = self._project_dir(member_id, project_id) / "docs" / f"{safe_filename(filename)}.txt"
        return path.read_text(encoding="utf-8") if path.is_file() else ""

    def build_context(self, member_id: str, project_id: str) -> str:
        """Bloc texte à injecter dans le system prompt pour « activer » un projet.

        Borné à CONTEXT_BUDGET : titres des documents + extraits, tronqués
        équitablement, pour ne pas faire exploser le contexte du LLM.
        """
        proj = self.get_project(member_id, project_id)
        if proj is None or not proj.documents:
            return ""
        budget_per_doc = max(400, self.CONTEXT_BUDGET // max(1, len(proj.documents)))
        chunks = [f"[Projet « {proj.name} » — documents fournis par l'utilisateur]"]
        for d in proj.documents:
            body = self.read_document(member_id, project_id, d.filename)
            if len(body) > budget_per_doc:
                body = body[:budget_per_doc].rstrip() + " […]"
            chunks.append(f"\n--- {d.filename} ---\n{body}")
        return "\n".join(chunks)

    def build_index(self, member_id: str, project_id: str) -> str:
        """Index LÉGER du projet (noms des fichiers seulement) pour le system prompt.

        Contrairement à build_context, n'injecte PAS le contenu : VindIA voit ce qui
        est disponible et lit À LA DEMANDE via ses outils (espace de travail à la
        Claude). Garde le contexte court même pour un gros dossier.
        """
        proj = self.get_project(member_id, project_id)
        if proj is None:
            return ""
        head = f"[Projet de référence actif : « {proj.name} »]"
        if not proj.documents:
            return head + "\nLe projet est vide. Tu peux y créer des fichiers (write_project_file)."
        files = "\n".join(f"- {d.filename}" for d in proj.documents)
        return (
            f"{head}\nFichiers disponibles (lis-les à la demande avec read_project_file, "
            f"n'invente jamais leur contenu) :\n{files}"
        )

    # -- persistance meta ---------------------------------------------------- #
    def _write_meta(self, member_id: str, proj: Project) -> None:
        meta = self._project_dir(member_id, proj.project_id) / "meta.json"
        meta.write_text(json.dumps(proj.as_dict(), ensure_ascii=False, indent=2), encoding="utf-8")

    @staticmethod
    def _read_meta(meta_path: Path) -> Project:
        data = json.loads(meta_path.read_text(encoding="utf-8"))
        docs = [Document(**d) for d in data.get("documents", [])]
        return Project(data["project_id"], data["name"], data.get("created_at", ""), docs)
